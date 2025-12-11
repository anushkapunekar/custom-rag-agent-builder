# backend/app/drive.py  (Corrected Final Multi-Document Version)

import os
import io
import json
from pathlib import Path
from typing import List

from fastapi import APIRouter, HTTPException, Request, Header
from fastapi.responses import JSONResponse

from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request as GoogleRequest
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload, HttpRequest

import fitz
import docx
from pptx import Presentation

from .auth import decode_token, get_user_by_id
from .utils import decrypt_json, ensure_user_dir
from .rag_utils import build_and_save_index   # MUST support full_text, doc_id, filename

router = APIRouter(prefix="/drive", tags=["drive"])

APP_DIR = Path(__file__).resolve().parent
STORAGE_BASE = APP_DIR / "storage"
STORAGE_BASE.mkdir(exist_ok=True)


# -----------------------------
# Google Drive service
# -----------------------------
def build_drive_service_from_creds(creds_json: dict):
    creds = Credentials(
        token=creds_json.get("token"),
        refresh_token=creds_json.get("refresh_token"),
        token_uri=creds_json.get("token_uri"),
        client_id=creds_json.get("client_id"),
        client_secret=creds_json.get("client_secret"),
        scopes=creds_json.get("scopes"),
    )

    if creds and creds.expired and creds.refresh_token:
        creds.refresh(GoogleRequest())

    return build("drive", "v3", credentials=creds, cache_discovery=False)


# -----------------------------
# TEXT CLEANING
# -----------------------------
def clean_text(t: str):
    if not t:
        return ""
    return (
        t.replace("\x00", "")
         .replace("\uf0b7", "")
         .replace("\u2022", "")
         .replace("\u25cf", "")
         .strip()
    )


# -----------------------------
# TEXT EXTRACTION
# -----------------------------
def extract_text_from_bytes(content: bytes, filename: str, mime: str):
    name = (filename or "").lower()

    # PDF
    if name.endswith(".pdf") or mime == "application/pdf":
        try:
            pdf = fitz.open(stream=content, filetype="pdf")
            text = ""
            for page in pdf:
                t = page.get_text("text")
                if not t.strip():
                    t = page.get_text("blocks")
                text += t + "\n"
            return clean_text(text)
        except Exception as e:
            print("PDF extract error:", e)
            return ""

    # DOCX
    if name.endswith(".docx"):
        try:
            d = docx.Document(io.BytesIO(content))
            return clean_text("\n".join(p.text for p in d.paragraphs))
        except Exception as e:
            print("DOCX extract error:", e)
            return ""

    # PPTX
    if name.endswith(".pptx"):
        try:
            prs = Presentation(io.BytesIO(content))
            txt = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        txt.append(shape.text)
            return clean_text("\n".join(txt))
        except Exception as e:
            print("PPTX extract error:", e)
            return ""

    # Fallback plain text
    try:
        return clean_text(content.decode("utf-8", errors="ignore"))
    except:
        return ""


# ===============================================================
#  /drive/download  — Upload + Extract + Per-Document Index Build
# ===============================================================
@router.post("/download")
async def download_and_index(request: Request, authorization: str = Header(None)):

    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(401, "Missing Authorization header")

    token = authorization.split(" ", 1)[1]
    payload = decode_token(token)
    user = get_user_by_id(payload["sub"])
    if not user:
        raise HTTPException(404, "User not found")

    try:
        body = await request.json()
    except:
        raw = await request.body()
        body = json.loads(raw.decode()) if raw else {}

    if "fileIds" not in body:
        raise HTTPException(400, "Body must contain fileIds")

    file_ids = body["fileIds"]

    gcreds = decrypt_json(user["creds"])
    service = build_drive_service_from_creds(gcreds)

    user_dir = ensure_user_dir(user["id"])
    docs_dir = user_dir / "docs"
    docs_dir.mkdir(exist_ok=True)

    uploaded = []
    total_chunks = 0

    for fid in file_ids:
        # -----------------------
        # first get metadata
        # -----------------------
        try:
            meta = service.files().get(fileId=fid, fields="id,name,mimeType").execute()
        except Exception as e:
            # skip this file but continue with others
            print(f"Failed to fetch metadata for {fid}: {e}")
            continue

        fname = meta.get("name") or f"{fid}"
        mime = meta.get("mimeType", "")

        # -----------------------
        # choose download method
        # -----------------------
        try:
            if mime == "application/vnd.google-apps.document":
                # export Google Doc as plain text
                request_obj: HttpRequest = service.files().export_media(fileId=fid, mimeType="text/plain")
            elif mime == "application/vnd.google-apps.spreadsheet":
                # export sheet as CSV
                request_obj = service.files().export_media(fileId=fid, mimeType="text/csv")
            elif mime == "application/vnd.google-apps.presentation":
                # export slides as plain text (best-effort)
                request_obj = service.files().export_media(fileId=fid, mimeType="text/plain")
            else:
                request_obj = service.files().get_media(fileId=fid)
        except Exception as e:
            print(f"Error preparing download for {fid} ({mime}): {e}")
            continue

        # -----------------------
        # download content
        # -----------------------
        fh = io.BytesIO()
        try:
            downloader = MediaIoBaseDownload(fh, request_obj)
            done = False
            while not done:
                _, done = downloader.next_chunk()
        except Exception as e:
            # some files (e.g. Google Forms or non-downloadable) may fail — skip
            print(f"Download failed for {fid}: {e}")
            continue

        content = fh.getvalue()

        # save raw file to user's docs dir (keep safe filename)
        safe_name = f"{fid}-{fname}"
        fpath = docs_dir / safe_name
        try:
            fpath.write_bytes(content)
        except Exception as e:
            print(f"Failed writing file {fpath}: {e}")

        # Extract text from bytes
        extracted = extract_text_from_bytes(content, fname, mime)

        # Build index (per-document) — use correct parameter names matching rag_utils
        try:
            added = build_and_save_index(
                user_id=user["id"],
                full_text=extracted,
                doc_id=fid,
                filename=fname
            )
        except Exception as e:
            print(f"Indexing failed for {fname}: {e}")
            # continue to next file rather than crash everything
            continue

        total_chunks += added
        uploaded.append({"id": fid, "name": fname})

    return {
        "status": "ok",
        "message": "Documents uploaded & indexed.",
        "files_uploaded": uploaded,
        "chunks_added": total_chunks
    }


# ===============================================================
#  /drive/list_docs — List stored user's documents
# ===============================================================
@router.get("/list_docs")
def list_docs(authorization: str = Header(None)):
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(401, "Missing Authorization header")

    token = authorization.split(" ", 1)[1]
    payload = decode_token(token)
    user = get_user_by_id(payload["sub"])

    if not user:
        raise HTTPException(404, "User not found")

    user_id = str(user["id"])
    user_dir = ensure_user_dir(user_id)
    docs_dir = user_dir / "docs"

    docs = []
    for f in docs_dir.iterdir():
        if not f.is_file():
            continue

        # filename format → "{docId}-{actualName}"
        parts = f.name.split("-", 1)
        if len(parts) == 2:
            docId, clean_name = parts
        else:
            docId = parts[0]
            clean_name = f.name

        docs.append({
            "id": docId,
            "name": clean_name,
            "path": str(f)
        })

    return {"docs": docs}
